import os
import pandas as pd
from docx import Document
from pptx import Presentation
import chardet
from fiber import *
import nltk
from nltk.tokenize import word_tokenize
from nltk.corpus import stopwords
from PyPDF2 import PdfReader  # New import for handling PDFs

# Download required NLTK data
nltk.download("punkt")
nltk.download("stopwords")
nltk.download('punkt_tab')

def indexing(LOCAL_CACHE_DIR):
    # Initialize the FiberDBMS instance
    dbms = FiberDBMS()
    temp_db_file = "temp_database.txt"
    dbms.load_or_create(temp_db_file)

    # Traverse the cache directory for all supported file types
    for root, _, files in os.walk(LOCAL_CACHE_DIR):
        for file in files:
            file_path = os.path.join(root, file)
            file_extension = os.path.splitext(file)[1].lower()

            try:
                # Process different file types and extract content
                if file_extension == ".txt":
                    with open(file_path, 'rb') as f:
                        raw_data = f.read()
                        encoding = chardet.detect(raw_data)['encoding']
                    with open(file_path, 'r', encoding=encoding) as f:
                        content = f.read()
                elif file_extension == ".docx":
                    doc = Document(file_path)
                    content = "\n".join([para.text for para in doc.paragraphs])
                elif file_extension == ".pptx":
                    presentation = Presentation(file_path)
                    content = "\n".join(
                        [slide.shapes.title.text if slide.shapes.title else '' 
                         + "\n".join([shape.text for shape in slide.shapes if hasattr(shape, "text")])
                         for slide in presentation.slides]
                    )
                elif file_extension in [".xls", ".xlsx", ".csv"]:
                    df = pd.read_excel(file_path) if "xls" in file_extension else pd.read_csv(file_path)
                    content = df.to_csv(index=False)
                elif file_extension == ".pdf":  # New condition for handling PDFs
                    reader = PdfReader(file_path)
                    content = "\n".join([page.extract_text() for page in reader.pages])
                else:
                    content = None  # Ignore unsupported formats

                # Add the file content to the database
                if content:  # Ensure content is not None
                    for i in content.split('\n'):
                        if i:
                            stop_words = set(stopwords.words('english'))
                            words = word_tokenize(i)
                            keywords = [word for word in words if word.lower() not in stop_words and word.isalpha()]  # Keep only relevant words

                            dbms.add_entry(name=file, content=i, tags=keywords)
                            print(i, keywords)
                print('The database has been indexed')
            except Exception as e:
                print(f"Failed to process {file}: {e}")

    # Save the database into a temp file
    dbms.save(temp_db_file)
