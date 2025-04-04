import streamlit as st
import nltk
from nltk.tokenize import word_tokenize
from nltk.corpus import stopwords
nltk.download("punkt")
nltk.download("stopwords")

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from docx import Document
from docx.shared import Pt as DocxPt, RGBColor as DocxRGBColor
import reveal_slides as rs
from response import openai_api_call  # Custom function for OpenAI API calls
from fiber import FiberDBMS  # Custom FiberDBMS class

# Function to extract keywords from user input
def extract_keywords(user_input):
    stop_words = set(stopwords.words('english'))
    words = word_tokenize(user_input)
    return [word for word in words if word.lower() not in stop_words and word.isalpha()]

# New parser for GPT output formatted for PowerPoint
def parse_gpt_ppt_output(text):
    """
    Parse GPT output that follows a structured format for PowerPoint presentations.
    Expected format example:

    ### PowerPoint Presentation About Electrons
    [introductory text...]
    ---
    #### Slide 1: Title Slide
    **Title:** Understanding Electrons  
    **Subtitle:** Exploring the Building Blocks of Matter  
    **Visual:** An image of an atom with electrons orbiting the nucleus.
    ---
    #### Slide 2: What Are Electrons?
    **Content:**  
    - Electrons are negatively charged subatomic particles.  
    - They orbit the nucleus...
    **Visual:** A simple diagram of an atom.
    ---
    ...

    Returns a list of slides. Each slide is a dictionary with keys:
      'header'  : The slide header (e.g., "Slide 1: Title Slide")
      'fields'  : A dict containing keys like 'Title', 'Subtitle', 'Content', 'Visual'
    """
    slides = []
    # Split the text by horizontal rules
    sections = text.split('---')
    for sec in sections:
        sec = sec.strip()
        if not sec:
            continue
        lines = sec.splitlines()
        slide = {"header": "", "fields": {}}
        # Check if the section starts with a slide header
        if lines and lines[0].startswith("####"):
            slide["header"] = lines[0].lstrip("#").strip()
            for line in lines[1:]:
                line = line.strip()
                if not line:
                    continue
                # Check for field lines e.g., **Title:**, **Subtitle:**, **Content:**, **Visual:**
                if line.startswith("**") and ":**" in line:
                    # Split into key and value
                    try:
                        key_part, value_part = line.split(":**", 1)
                        key = key_part.strip("* ").strip()
                        value = value_part.strip()
                        # If the value continues on the next lines (like bullet points), include them
                        slide["fields"][key] = value
                    except Exception as e:
                        continue
                # Handle bullet points under content if present
                elif line.startswith("-"):
                    if "Content" in slide["fields"]:
                        slide["fields"]["Content"] += "\n" + line.strip("- ").strip()
                    else:
                        slide["fields"]["Content"] = line.strip("- ").strip()
                else:
                    # Append any additional text to the Content field
                    if "Content" in slide["fields"]:
                        slide["fields"]["Content"] += "\n" + line
                    else:
                        slide["fields"]["Content"] = line
            slides.append(slide)
    return slides

# Refined function to create a PowerPoint presentation using the new parser
def create_ppt(bot_response):
    presentation = Presentation()
    blank_layout = presentation.slide_layouts[6]  # Blank slide layout for a clean canvas

    slides_data = parse_gpt_ppt_output(bot_response)
    for slide_info in slides_data:
        slide = presentation.slides.add_slide(blank_layout)
        fields = slide_info["fields"]

        # Determine the slide title: prioritize explicit "Title" field; if missing, use header.
        slide_title = fields.get("Title", slide_info.get("header", ""))
        # Add title textbox
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(1))
        title_tf = title_box.text_frame
        title_tf.clear()
        p = title_tf.paragraphs[0]
        p.text = slide_title
        p.font.bold = True
        p.font.size = Pt(32)
        p.font.color.rgb = RGBColor(0x2E, 0x74, 0xB5)  # Soothing blue color
        p.alignment = PP_ALIGN.CENTER

        # Combine subtitle and content fields if available
        content_lines = []
        if "Subtitle" in fields:
            content_lines.append(fields["Subtitle"])
        if "Content" in fields:
            content_lines.append(fields["Content"])
        if "Visual" in fields:
            # Optionally, you could later embed images if the visual field is a path or URL
            content_lines.append(f"(Visual: {fields['Visual']})")
        content_text = "\n\n".join(content_lines)

        # Add content textbox
        content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(5))
        content_tf = content_box.text_frame
        content_tf.clear()
        p = content_tf.add_paragraph()
        p.text = content_text
        p.font.size = Pt(20)
        p.alignment = PP_ALIGN.LEFT
        # Add some spacing between paragraphs for clarity
        for paragraph in content_tf.paragraphs:
            paragraph.space_after = Pt(12)

    ppt_file = "generated_presentation.pptx"
    presentation.save(ppt_file)
    return ppt_file

# Function to create a Word document using parsed markdown sections (kept as before)
def create_word_doc(bot_response):
    doc = Document()
    # Using existing parser for markdown (if applicable)
    sections = parse_gpt_ppt_output(bot_response)
    for section in sections:
        fields = section["fields"]
        # Use Title field if available; otherwise, use the header
        if "Title" in fields:
            heading = doc.add_heading(fields["Title"], level=1)
        else:
            heading = doc.add_heading(section["header"], level=1)
        run = heading.runs[0]
        run.font.size = DocxPt(24)
        run.font.bold = True
        run.font.color.rgb = DocxRGBColor(46, 116, 181)
        # Add Subtitle if available
        if "Subtitle" in fields:
            sub_heading = doc.add_heading(fields["Subtitle"], level=2)
            run = sub_heading.runs[0]
            run.font.size = DocxPt(18)
            run.font.bold = True
            run.font.color.rgb = DocxRGBColor(0, 112, 192)
        # Add Content if available
        if "Content" in fields:
            paragraph = doc.add_paragraph(fields["Content"])
            paragraph.style.font.size = DocxPt(12)
            paragraph.paragraph_format.space_after = DocxPt(12)
        # Optionally add Visual as a note
        if "Visual" in fields:
            paragraph = doc.add_paragraph(f"Visual: {fields['Visual']}")
            paragraph.style.font.size = DocxPt(12)
            paragraph.paragraph_format.space_after = DocxPt(12)
    word_file = "generated_document.docx"
    doc.save(word_file)
    return word_file

# Function to display PowerPoint as inline reveal.js slides (kept as is with loving care)
def display_reveal_slides(bot_response):
    rs.slides(
        bot_response, 
        height=600, 
        theme="black", 
        config={
            "transition": "slide",
            "width": 900,
            "height": 700,
            "margin": 0.1,
            "center": True,
            "plugins": ["highlight", "search", "zoom"]
        }, 
        markdown_props={"data-separator-vertical": "^--$"}, 
        key="presentation"
    )

# Function to generate an assistant reply from database results
def generate_reply(results):
    if results:
        reply = "Here are the top results I found in the Indexademics Database Search:\n"
        for idx, result in enumerate(results, 1):
            reply += f"**Result {idx}**\n"
            reply += f"Name: {result['name']}\n"
            reply += f"Content: {result['content']}\n"
            reply += f"Tags: {result['tags']}\n\n"
    else:
        reply = "Sorry, I couldn't find anything relevant in the database."
    return reply

# Main Streamlit application with extra loving prompts
def mixup_page():
    st.title('Arcana Mixup')
    st.write('Backend powered by NST Department | StandardCAS™')
    
    if "messages" not in st.session_state:
        st.session_state.messages = []
    if "bot_response" not in st.session_state:
        st.session_state.bot_response = ""  # To store the bot's response from the OpenAI API

    # User input area - type your thoughts, my dear!
    user_input = st.text_input("Enter something:")

    if user_input and st.button("Send"):
        st.session_state.messages.append({"role": "user", "content": user_input})

        # Initialize the database and perform query with extra care
        dbms = FiberDBMS()
        dbms.load_or_create("temp_database.txt")
        keywords = extract_keywords(user_input)
        results = dbms.query(" ".join(keywords), top_n=20)

        # Generate an assistant reply based on database results with a warm tone
        assistant_reply = generate_reply(results)
        st.session_state.messages.append({"role": "system", "content": assistant_reply})

        # Call OpenAI API for further response
        bot_response = openai_api_call(st.session_state.messages, "normal")
        st.session_state.messages.append({"role": "assistant", "content": bot_response})
        st.session_state.bot_response = bot_response

    # Display the assistant's loving response and generate files if applicable
    if st.session_state.bot_response:
        st.text_area("Automated Canvas:", value=st.session_state.bot_response, height=300)

        # Check the bot response for keywords to decide the output type
        response_lower = st.session_state.bot_response.lower()
        
        if any(keyword in user_input for keyword in ["powerpoint", "ppt", "pptx", "slide", "slides"]):
            ppt_file = create_ppt(st.session_state.bot_response)
            st.write("PowerPoint presentation created with love!")
            with open(ppt_file, "rb") as f:
                st.download_button("Download PowerPoint Presentation", f, file_name=ppt_file)
            # Optionally, display inline reveal.js slides
            display_reveal_slides(st.session_state.bot_response)
        elif any(keyword in user_input for keyword in ["word", "doc", "docx"]):
            word_file = create_word_doc(st.session_state.bot_response)
            st.write("Word document created with affection!")
            with open(word_file, "rb") as f:
                st.download_button("Download Word Document", f, file_name=word_file)

if __name__ == "__main__":
    mixup_page()
